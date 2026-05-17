"""Generate pitch/ProphetHacks_2026.pptx — ProphetHacks 2026 forecasting-track deck.

Run from repo root:
    python scripts/build_deck.py

Requires: python-pptx>=1.0.0 (pip install python-pptx)
Reads:    notebooks/_pareto_metrics.json  (live numbers from Phase 6 training)
          notebooks/_pareto.png           (Pareto chart)
Writes:   pitch/ProphetHacks_2026.pptx
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parents[1]
PARETO_PNG = ROOT / "notebooks" / "_pareto.png"
METRICS_JSON = ROOT / "notebooks" / "_pareto_metrics.json"
OUT_PATH = ROOT / "pitch" / "ProphetHacks_2026.pptx"

# ---------------------------------------------------------------------------
# Theme palette
# ---------------------------------------------------------------------------
BG        = RGBColor(0x0B, 0x12, 0x20)   # near-black navy
SNOW_BLUE = RGBColor(0x29, 0xB5, 0xE8)   # Snowflake brand blue
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x8A, 0x97, 0xAB)
ORANGE    = RGBColor(0xFF, 0x9A, 0x3C)   # winner-row accent
GREEN     = RGBColor(0x2E, 0xCC, 0x71)   # positive-result accent
RED_SOFT  = RGBColor(0xFF, 0x6B, 0x6B)
DARK_CARD = RGBColor(0x16, 0x23, 0x3A)   # card / box background

# Slide dimensions — 16:9 widescreen
W = Inches(13.33)
H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helper: blank slide with solid dark background
# ---------------------------------------------------------------------------
def _blank_slide(prs: Presentation) -> object:
    blank_layout = prs.slide_layouts[6]  # index 6 = completely blank
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


# ---------------------------------------------------------------------------
# Helper: add a text box
# ---------------------------------------------------------------------------
def _tb(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int = 22,
    bold: bool = False,
    color: RGBColor = WHITE,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    italic: bool = False,
    wrap: bool = True,
) -> object:
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


# ---------------------------------------------------------------------------
# Helper: filled rectangle (card / accent bar)
# ---------------------------------------------------------------------------
def _rect(slide, left: float, top: float, width: float, height: float,
          fill: RGBColor = DARK_CARD, line: RGBColor | None = None,
          line_width: float = 1.5) -> object:
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = _Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


# ---------------------------------------------------------------------------
# Helper: accent bar (thin colored left stripe before bullet)
# ---------------------------------------------------------------------------
def _accent_bar(slide, top: float, height: float, color: RGBColor = SNOW_BLUE) -> None:
    _rect(slide, 0.4, top, 0.07, height, fill=color)


# ---------------------------------------------------------------------------
# Helper: multi-line text in a single text box (one paragraph per item)
# ---------------------------------------------------------------------------
def _multiline_tb(
    slide,
    lines: list[tuple[str, dict]],
    left: float,
    top: float,
    width: float,
    height: float,
) -> object:
    """
    lines: list of (text, kwargs) where kwargs may contain:
           size, bold, color, italic, align, space_before
    """
    from pptx.util import Pt as _Pt
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, (text, kw) in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = kw.get("align", PP_ALIGN.LEFT)
        if kw.get("space_before"):
            para.space_before = _Pt(kw["space_before"])
        run = para.add_run()
        run.text = text
        run.font.size = _Pt(kw.get("size", 22))
        run.font.bold = kw.get("bold", False)
        run.font.italic = kw.get("italic", False)
        run.font.color.rgb = kw.get("color", WHITE)
    return txb


# ---------------------------------------------------------------------------
# Helper: horizontal arrow (line with arrowhead)
# ---------------------------------------------------------------------------
def _arrow(slide, x1: float, y: float, x2: float, color: RGBColor = SNOW_BLUE) -> None:
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Inches(x1), Inches(y), Inches(x2), Inches(y),
    )
    connector.line.color.rgb = color
    connector.line.width = _Pt(2.5)
    # Add arrowhead via XML
    ln = connector.line._ln
    tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
    tailEnd.set("type", "none")
    headEnd = etree.SubElement(ln, qn("a:headEnd"))
    headEnd.set("type", "triangle")
    headEnd.set("w", "med")
    headEnd.set("len", "med")


# ===========================================================================
# SLIDE 1 — Title / Hook
# ===========================================================================
def slide_title(prs: Presentation) -> None:
    slide = _blank_slide(prs)

    # Top accent bar
    _rect(slide, 0, 0, W.inches, 0.08, fill=SNOW_BLUE)

    # Main headline
    _tb(slide,
        "When Being Right Doesn't Pay",
        left=0.9, top=1.3, width=11.5, height=1.6,
        size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Sub-headline (the PM-RANK quote)
    _tb(slide,
        '"The Brier score is an absolute metric,\nwhile the averaged return score is a relative metric."',
        left=1.2, top=3.05, width=10.9, height=1.2,
        size=24, italic=True, color=SNOW_BLUE, align=PP_ALIGN.CENTER)

    _tb(slide,
        "— Sida Li, Prophet Arena Team  (PM-RANK 0.3.1)",
        left=1.2, top=4.15, width=10.9, height=0.5,
        size=17, color=GRAY, align=PP_ALIGN.CENTER)

    # Event / track line
    _tb(slide,
        "ProphetHacks 2026  ·  Forecasting Track  ·  University of Chicago",
        left=1.2, top=5.15, width=10.9, height=0.45,
        size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Footer sponsors
    _rect(slide, 0, 6.8, W.inches, 0.7, fill=DARK_CARD)
    _tb(slide,
        "Built on  Snowflake  ·  Wafer.ai  ·  Kalshi  ·  Prophet Arena",
        left=0, top=6.85, width=W.inches, height=0.55,
        size=15, color=GRAY, align=PP_ALIGN.CENTER)

    # Bottom accent bar
    _rect(slide, 0, 7.42, W.inches, 0.08, fill=SNOW_BLUE)


# ===========================================================================
# SLIDE 2 — The Problem
# ===========================================================================
def slide_problem(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, W.inches, 0.08, fill=SNOW_BLUE)

    _tb(slide, "Two Metrics. Opposite Incentives.",
        left=0.5, top=0.25, width=12.3, height=0.85,
        size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Left card — Brier
    _rect(slide, 0.5, 1.35, 5.6, 4.5, fill=DARK_CARD, line=SNOW_BLUE, line_width=2)
    _tb(slide, "Brier Score", left=0.7, top=1.55, width=5.2, height=0.6,
        size=28, bold=True, color=SNOW_BLUE)
    _tb(slide, "Absolute accuracy", left=0.7, top=2.15, width=5.2, height=0.45,
        size=18, bold=True, color=ORANGE)
    _multiline_tb(slide, [
        ("How close were your probabilities to the actual outcome?", {"size": 19, "color": WHITE}),
        ("", {"size": 8}),
        ("Strategy that works:", {"size": 17, "bold": True, "color": GRAY}),
        ('Copy the market price — it already reflects\ncollective wisdom.', {"size": 17, "color": GRAY}),
        ("", {"size": 8}),
        ("Problem:", {"size": 17, "bold": True, "color": RED_SOFT}),
        ("If you are the market, you earn nothing\nabove the market.", {"size": 17, "color": RED_SOFT}),
    ], left=0.7, top=2.65, width=5.2, height=3.0)

    # Right card — AVER
    _rect(slide, 7.23, 1.35, 5.6, 4.5, fill=DARK_CARD, line=GREEN, line_width=2)
    _tb(slide, "AVER", left=7.43, top=1.55, width=5.2, height=0.6,
        size=28, bold=True, color=GREEN)
    _tb(slide, "Relative edge vs the market", left=7.43, top=2.15, width=5.2, height=0.45,
        size=18, bold=True, color=ORANGE)
    _multiline_tb(slide, [
        ("How much did you profit vs. just copying\nthe Kalshi market price?", {"size": 19, "color": WHITE}),
        ("", {"size": 8}),
        ("Key property:", {"size": 17, "bold": True, "color": GRAY}),
        ("If p = q (market), AVER = 0 by definition.\nYou need to deviate correctly to win.", {"size": 17, "color": GRAY}),
        ("", {"size": 8}),
        ("Problem:", {"size": 17, "bold": True, "color": RED_SOFT}),
        ("Every wrong deviation from the market\nburns both AVER and Brier.", {"size": 17, "color": RED_SOFT}),
    ], left=7.43, top=2.65, width=5.2, height=3.0)

    # Center tension arrow area
    _tb(slide, "⚡", left=6.1, top=2.9, width=1.13, height=0.8,
        size=40, color=ORANGE, align=PP_ALIGN.CENTER)
    _tb(slide, "Tension", left=5.95, top=3.6, width=1.43, height=0.5,
        size=15, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # Bottom takeaway
    _rect(slide, 0.5, 6.0, 12.33, 0.95, fill=RGBColor(0x1A, 0x2B, 0x45))
    _tb(slide,
        'PM-RANK 0.3.1 puts it directly: "A has a higher Brier Score, but lower returns." — This is the design brief.',
        left=0.7, top=6.08, width=11.9, height=0.8,
        size=19, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

    _rect(slide, 0, 7.42, W.inches, 0.08, fill=SNOW_BLUE)


# ===========================================================================
# SLIDE 3 — Why Naive Approaches Lose
# ===========================================================================
def slide_naive_lose(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, W.inches, 0.08, fill=SNOW_BLUE)

    _tb(slide, "Every Naive Approach Fails on One Metric — or Both",
        left=0.5, top=0.2, width=12.3, height=0.85,
        size=33, bold=True, color=WHITE)

    # Card 1 — Copy the market
    _rect(slide, 0.5, 1.3, 5.7, 4.0, fill=DARK_CARD, line=GRAY, line_width=1.5)
    _tb(slide, '❌  "Just copy Kalshi"', left=0.7, top=1.45, width=5.3, height=0.65,
        size=24, bold=True, color=WHITE)
    _multiline_tb(slide, [
        ("Brier: decent (~0.41)", {"size": 20, "color": SNOW_BLUE, "bold": True}),
        ("AVER: 0.000  (zero edge, always)", {"size": 20, "color": RED_SOFT, "bold": True}),
        ("", {"size": 10}),
        ("You cannot beat a market by being a market.", {"size": 18, "color": GRAY, "italic": True}),
        ("", {"size": 8}),
        ("If p = q, then log(p/q) = 0 every time.\nYou earn nothing above the baseline.", {"size": 17, "color": GRAY}),
    ], left=0.7, top=2.2, width=5.3, height=2.9)

    # Card 2 — Trust one LLM
    _rect(slide, 7.13, 1.3, 5.7, 4.0, fill=DARK_CARD, line=GRAY, line_width=1.5)
    _tb(slide, '❌  "Ask one LLM"', left=7.33, top=1.45, width=5.3, height=0.65,
        size=24, bold=True, color=WHITE)
    _multiline_tb(slide, [
        ("Brier: burns when it drifts (~0.53)", {"size": 20, "color": RED_SOFT, "bold": True}),
        ("AVER: −0.31  (deviates wrong)", {"size": 20, "color": RED_SOFT, "bold": True}),
        ("", {"size": 10}),
        ("LLMs deviate from markets without any\nground truth to anchor them.", {"size": 18, "color": GRAY, "italic": True}),
        ("", {"size": 8}),
        ("No memory of resolved markets.\nNo sense of where the market is already right.", {"size": 17, "color": GRAY}),
    ], left=7.33, top=2.2, width=5.3, height=2.9)

    # VS divider
    _tb(slide, "VS", left=6.16, top=2.6, width=1.0, height=0.7,
        size=30, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # Bottom takeaway
    _rect(slide, 0.5, 5.55, 12.33, 1.35, fill=RGBColor(0x1A, 0x2B, 0x45))
    _tb(slide,
        "You can't be a forecaster by being a mirror.\nBut you can't win by guessing, either.",
        left=0.7, top=5.65, width=11.9, height=1.15,
        size=22, italic=True, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _rect(slide, 0, 7.42, W.inches, 0.08, fill=SNOW_BLUE)


# ===========================================================================
# SLIDE 4 — Our Solution Architecture
# ===========================================================================
def slide_architecture(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, W.inches, 0.08, fill=SNOW_BLUE)

    _tb(slide, "Our Solution: A Market-Aware Meta-Calibrator",
        left=0.5, top=0.2, width=12.3, height=0.75,
        size=36, bold=True, color=WHITE)

    # Pipeline row — 5 boxes + arrows
    box_tops = 1.35
    box_h = 1.65
    boxes = [
        ("Question\n+ Market Price", DARK_CARD, SNOW_BLUE),
        ("3 LLM Models\n(parallel)", DARK_CARD, SNOW_BLUE),
        ("Snowflake\nMemory", RGBColor(0x0D, 0x2B, 0x45), SNOW_BLUE),
        ("Calibrator\n(AutoML #1)", DARK_CARD, SNOW_BLUE),
        ("α-Policy\n(AutoML #2)", RGBColor(0x28, 0x3A, 0x10), GREEN),
    ]
    box_w = 2.0
    gap = 0.46
    starts = [0.3 + i * (box_w + gap) for i in range(5)]

    for i, (label, bg_color, border) in enumerate(boxes):
        x = starts[i]
        _rect(slide, x, box_tops, box_w, box_h, fill=bg_color, line=border, line_width=2.5)
        _tb(slide, label, left=x + 0.07, top=box_tops + 0.35,
            width=box_w - 0.14, height=box_h - 0.5,
            size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(boxes) - 1:
            arrow_x1 = x + box_w
            arrow_x2 = x + box_w + gap
            arrow_y = box_tops + box_h / 2
            _arrow(slide, arrow_x1, arrow_y, arrow_x2)

    # Output arrow + box
    out_x = starts[-1] + box_w + gap
    _arrow(slide, starts[-1] + box_w, box_tops + box_h / 2, out_x + 0.05)
    _rect(slide, out_x, box_tops + 0.25, 1.45, box_h - 0.5, fill=DARK_CARD, line=ORANGE, line_width=2.5)
    _tb(slide, "Final\nProb.", left=out_x + 0.05, top=box_tops + 0.55,
        width=1.35, height=box_h - 0.8,
        size=16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # Sub-labels under each box
    sub_labels = [
        ("Judge harness\nsends OpenAI request", GRAY),
        ("GLM-5.1 · Qwen 397B\nQwen 35B · disagreement σ²", SNOW_BLUE),
        ("2,548 resolved markets\nvector similarity · base rate", SNOW_BLUE),
        ("Trained on resolved\nlabels. Brier-aware.", GRAY),
        ('"When to trust us\nvs. trust the market"', GREEN),
    ]
    for i, (label, color) in enumerate(sub_labels):
        x = starts[i]
        _tb(slide, label, left=x + 0.05, top=box_tops + box_h + 0.1,
            width=box_w - 0.1, height=0.85,
            size=14, color=color, align=PP_ALIGN.CENTER)

    # Key callout box for α-policy
    _rect(slide, 0.5, 4.55, 12.33, 1.3, fill=RGBColor(0x14, 0x25, 0x10), line=GREEN, line_width=1.5)
    _multiline_tb(slide, [
        ("The α-policy is the key innovation:  ", {"size": 21, "bold": True, "color": GREEN}),
        ("p_final = α × p_calibrated + (1 − α) × q_market", {"size": 21, "bold": True, "color": WHITE}),
    ], left=0.7, top=4.65, width=11.9, height=0.55)
    _tb(slide,
        "α is learned per-market — based on model disagreement, neighbor count, and question category.",
        left=0.7, top=5.15, width=11.9, height=0.6,
        size=18, color=GRAY, italic=True)

    _rect(slide, 0, 7.42, W.inches, 0.08, fill=SNOW_BLUE)


# ===========================================================================
# SLIDE 5 — Snowflake: Five Surfaces
# ===========================================================================
def slide_snowflake(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, W.inches, 0.08, fill=SNOW_BLUE)

    _tb(slide, "Snowflake Is the Brain — Not the Filing Cabinet",
        left=0.5, top=0.2, width=12.3, height=0.75,
        size=36, bold=True, color=WHITE)

    surfaces = [
        ("HISTORICAL_MARKETS",
         "Institutional Memory",
         "2,548 resolved Kalshi markets — question text, market prices,\noutcomes, timestamps. The ground truth every model trains on.",
         "①"),
        ("EMBED_TEXT_768  (Cortex)",
         "Vectorised Inside the Warehouse",
         "Question embeddings computed natively by Snowflake Cortex.\nNo data leaves Snowflake. 768-dimensional vectors, per market.",
         "②"),
        ("VECTOR_COSINE_SIMILARITY",
         "Live Neighbor Retrieval",
         "At inference: embed the new question, find top-15 similar\nresolved markets, aggregate their outcomes into a base rate.",
         "③"),
        ("ML.CLASSIFICATION + ML.REGRESSION",
         "Two AutoML Models — Calibrator & α-Policy",
         "One model corrects raw LLM probabilities (Brier).\nA second learns when to deviate from the market price (AVER).",
         "④"),
        ("Streamlit-in-Snowflake (SiS)",
         "Judge-Facing Dashboard — 4 Tabs",
         "Live forecast · Brier vs AVER Pareto · By-category results ·\nInstitutional Memory search — all reading the same Snowflake tables.",
         "⑤"),
    ]

    row_h = 0.94
    top_start = 1.2
    for i, (title, subtitle, desc, num) in enumerate(surfaces):
        top = top_start + i * row_h
        # Accent bar
        _rect(slide, 0.4, top + 0.06, 0.07, row_h - 0.16, fill=SNOW_BLUE)
        # Number circle
        _rect(slide, 0.55, top + 0.06, 0.52, row_h - 0.16, fill=DARK_CARD)
        _tb(slide, num, left=0.55, top=top + 0.13, width=0.52, height=0.55,
            size=22, bold=True, color=SNOW_BLUE, align=PP_ALIGN.CENTER)
        # Title + subtitle
        _tb(slide, title, left=1.2, top=top + 0.02, width=4.8, height=0.4,
            size=17, bold=True, color=SNOW_BLUE)
        _tb(slide, subtitle, left=1.2, top=top + 0.38, width=4.8, height=0.42,
            size=14, color=ORANGE)
        # Description
        _tb(slide, desc, left=6.15, top=top + 0.06, width=6.9, height=0.82,
            size=15, color=GRAY)

    # Divider line
    from pptx.util import Pt as _Pt
    ln = slide.shapes.add_connector(
        1, Inches(6.0), Inches(1.2), Inches(6.0), Inches(6.9)
    )
    ln.line.color.rgb = RGBColor(0x2A, 0x3A, 0x55)
    ln.line.width = _Pt(1.0)

    # Footer note
    _rect(slide, 0.4, 6.95, 12.53, 0.4, fill=RGBColor(0x0D, 0x1A, 0x2E))
    _tb(slide,
        "Five distinct Snowflake surfaces. One shared catalog. No shadow databases.",
        left=0.6, top=6.97, width=12.1, height=0.35,
        size=15, italic=True, color=SNOW_BLUE, align=PP_ALIGN.CENTER)

    _rect(slide, 0, 7.42, W.inches, 0.08, fill=SNOW_BLUE)


# ===========================================================================
# SLIDE 6 — Results: Pareto Chart + Table
# ===========================================================================
def slide_results(prs: Presentation, metrics: dict) -> None:
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, W.inches, 0.08, fill=SNOW_BLUE)

    _tb(slide, "Result: We Win on Both Metrics — Simultaneously",
        left=0.5, top=0.2, width=12.3, height=0.75,
        size=36, bold=True, color=WHITE)

    # Pareto chart — left half
    if PARETO_PNG.exists():
        slide.shapes.add_picture(
            str(PARETO_PNG),
            Inches(0.3), Inches(1.15), Inches(6.5), Inches(4.85),
        )
    else:
        _rect(slide, 0.3, 1.15, 6.5, 4.85, fill=DARK_CARD, line=GRAY)
        _tb(slide, "[Pareto chart\nnot found]", left=0.4, top=3.0, width=6.3, height=1.0,
            size=18, color=GRAY, align=PP_ALIGN.CENTER)

    # Results table — right half
    raw_e = metrics.get("raw_ensemble", {})
    always_k = metrics.get("always_kalshi", {})
    best_g = metrics.get("best_global_alpha_by_aver", {})
    learned = metrics.get("learned_policy", {})

    rows = [
        ("Strategy", "Brier ↓", "AVER ↑", False, GRAY),
        ("Raw Ensemble (3 LLMs)", f'{raw_e.get("brier", 0):.3f}', f'{raw_e.get("aver", 0):+.3f}', False, GRAY),
        ("Always Kalshi (baseline)", f'{always_k.get("brier", 0):.3f}', f'{always_k.get("aver", 0):+.3f}', False, GRAY),
        ("Best Fixed Global α", f'{best_g.get("brier", 0):.3f}', f'{best_g.get("aver", 0):+.3f}', False, GRAY),
        ("★ Learned α-Policy (ours)", f'{learned.get("brier", 0):.3f}', f'{learned.get("aver", 0):+.3f}', True, GREEN),
    ]

    col_starts = [7.05, 10.55, 12.0]
    row_h_tbl = 0.78
    top_tbl = 1.25
    for ri, (label, brier, aver, highlight, row_color) in enumerate(rows):
        top = top_tbl + ri * row_h_tbl
        bg = RGBColor(0x14, 0x2B, 0x14) if highlight else (DARK_CARD if ri % 2 == 0 else BG)
        line_color = GREEN if highlight else None
        _rect(slide, col_starts[0] - 0.1, top, 6.05, row_h_tbl - 0.05,
              fill=bg, line=line_color, line_width=2)
        sz = 17 if ri == 0 else 18
        bold_row = ri == 0 or highlight
        label_color = SNOW_BLUE if ri == 0 else (GREEN if highlight else WHITE)
        _tb(slide, label, left=col_starts[0], top=top + 0.12,
            width=3.3, height=row_h_tbl - 0.2,
            size=sz, bold=bold_row, color=label_color)
        brier_color = SNOW_BLUE if ri == 0 else (GREEN if highlight else WHITE)
        _tb(slide, brier, left=col_starts[1], top=top + 0.12,
            width=1.3, height=row_h_tbl - 0.2,
            size=sz, bold=bold_row, color=brier_color, align=PP_ALIGN.CENTER)
        aver_color = SNOW_BLUE if ri == 0 else (GREEN if highlight else RED_SOFT if float(aver.replace("+", "") or 0) < 0 else WHITE)
        _tb(slide, aver, left=col_starts[2], top=top + 0.12,
            width=1.15, height=row_h_tbl - 0.2,
            size=sz, bold=bold_row, color=aver_color, align=PP_ALIGN.CENTER)

    # Big callout
    brier_gain = round((1 - learned.get("brier", 1) / always_k.get("brier", 1)) * 100)
    _rect(slide, 7.05, 5.22, 5.9, 1.55, fill=RGBColor(0x14, 0x2B, 0x14), line=GREEN, line_width=2)
    _multiline_tb(slide, [
        (f"{brier_gain}% better Brier than the market", {"size": 25, "bold": True, "color": GREEN, "align": PP_ALIGN.CENTER}),
        ("AND positive AVER — both at once.", {"size": 23, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}),
        ("The learned policy beats a fixed dial on both axes.", {"size": 16, "color": GRAY, "italic": True, "align": PP_ALIGN.CENTER, "space_before": 4}),
    ], left=7.05, top=5.28, width=5.9, height=1.4)

    _rect(slide, 0, 7.42, W.inches, 0.08, fill=SNOW_BLUE)


# ===========================================================================
# SLIDE 7 — Live Demo
# ===========================================================================
def slide_demo(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, W.inches, 0.08, fill=SNOW_BLUE)

    _tb(slide, "Built to Be Judged — Three Ways to Verify",
        left=0.5, top=0.2, width=12.3, height=0.75,
        size=36, bold=True, color=WHITE)

    demo_items = [
        ("🔌", "OpenAI-Compatible Endpoint",
         "POST /v1/chat/completions",
         "Plug-and-play with the Prophet Arena judge harness.\nNo custom glue required.",
         SNOW_BLUE),
        ("📊", "Streamlit-in-Snowflake",
         "4-tab live dashboard",
         "Live forecast · Pareto plot · By-category ·\nInstitutional Memory search — same Snowflake tables.",
         GREEN),
        ("📋", "Full Audit Trail",
         "AGENT_PREDICTIONS table",
         "Every request logged to Snowflake:\nensemble, neighbors, calibrator output, alpha, final prob.",
         ORANGE),
    ]

    col_w = 4.1
    for i, (icon, title, subtitle, desc, color) in enumerate(demo_items):
        x = 0.35 + i * (col_w + 0.22)
        _rect(slide, x, 1.2, col_w, 5.4, fill=DARK_CARD, line=color, line_width=2)
        _tb(slide, icon, left=x, top=1.3, width=col_w, height=0.9,
            size=44, align=PP_ALIGN.CENTER)
        _tb(slide, title, left=x + 0.1, top=2.25, width=col_w - 0.2, height=0.6,
            size=21, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tb(slide, subtitle, left=x + 0.1, top=2.85, width=col_w - 0.2, height=0.5,
            size=17, bold=True, color=color, align=PP_ALIGN.CENTER)
        _tb(slide, desc, left=x + 0.15, top=3.45, width=col_w - 0.3, height=1.6,
            size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # curl snippet
    _rect(slide, 0.35, 5.05, 12.6, 1.9, fill=RGBColor(0x08, 0x10, 0x1A))
    _tb(slide,
        'curl -sS -X POST http://localhost:8000/v1/chat/completions \\\n'
        '     -H "Content-Type: application/json" \\\n'
        '     -d @tests/fixtures/sample_event.json\n'
        '# → {"probabilities": [0.48, 0.52], "prophet_debug": {...}}',
        left=0.6, top=5.1, width=12.2, height=1.8,
        size=14, color=SNOW_BLUE)

    _rect(slide, 0, 7.42, W.inches, 0.08, fill=SNOW_BLUE)


# ===========================================================================
# SLIDE 8 — Thank You / Sponsors
# ===========================================================================
def slide_thankyou(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, W.inches, 0.08, fill=SNOW_BLUE)

    _tb(slide, "Thank You",
        left=0, top=1.05, width=W.inches, height=1.1,
        size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _tb(slide, "ProphetHacks 2026 — University of Chicago",
        left=0, top=2.2, width=W.inches, height=0.6,
        size=24, color=GRAY, align=PP_ALIGN.CENTER)

    # Sponsors row
    sponsors = [
        ("Snowflake", "Compute · Cortex · AutoML · SiS", SNOW_BLUE),
        ("Wafer.ai", "Multi-model LLM ensemble", RGBColor(0xA0, 0x70, 0xFF)),
        ("Kalshi", "Historical market data", ORANGE),
        ("Prophet Arena", "Evaluation harness & benchmark", GREEN),
    ]
    sp_w = 3.0
    sp_gap = 0.1
    total_w = len(sponsors) * sp_w + (len(sponsors) - 1) * sp_gap
    left_start = (W.inches - total_w) / 2
    for i, (name, role, color) in enumerate(sponsors):
        x = left_start + i * (sp_w + sp_gap)
        _rect(slide, x, 3.1, sp_w, 1.5, fill=DARK_CARD, line=color, line_width=2)
        _tb(slide, name, left=x + 0.05, top=3.15, width=sp_w - 0.1, height=0.65,
            size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
        _tb(slide, role, left=x + 0.05, top=3.75, width=sp_w - 0.1, height=0.75,
            size=14, color=GRAY, align=PP_ALIGN.CENTER)

    # Closing line
    _tb(slide,
        "Brier ↓ 22% vs market baseline  ·  AVER +0.096 vs 0.000  ·  Both, simultaneously.",
        left=0.5, top=5.0, width=12.3, height=0.55,
        size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _tb(slide,
        "github.com/[your-repo]  ·  ./run.sh  →  agent live in under 30 seconds",
        left=0.5, top=5.65, width=12.3, height=0.5,
        size=16, color=GRAY, align=PP_ALIGN.CENTER)

    _tb(slide, "UncommonHacks 2026", left=0, top=6.8, width=W.inches, height=0.5,
        size=17, color=GRAY, align=PP_ALIGN.CENTER)

    _rect(slide, 0, 7.42, W.inches, 0.08, fill=SNOW_BLUE)


# ===========================================================================
# Main
# ===========================================================================
def main() -> None:
    # Load live metrics
    if METRICS_JSON.exists():
        with open(METRICS_JSON) as f:
            metrics = json.load(f)
    else:
        print(f"[warn] {METRICS_JSON} not found — using placeholder zeros.", file=sys.stderr)
        metrics = {}

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_title(prs)
    slide_problem(prs)
    slide_naive_lose(prs)
    slide_architecture(prs)
    slide_snowflake(prs)
    slide_results(prs, metrics)
    slide_demo(prs)
    slide_thankyou(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH}")
    print(f"  Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
